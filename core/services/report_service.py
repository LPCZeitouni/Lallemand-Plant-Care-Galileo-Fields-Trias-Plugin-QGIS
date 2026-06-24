# -*- coding: utf-8 -*-
"""
LallemandGeostatFieldTrialTreatments

Report service for Word and PowerPoint generation.
"""

import os.path
import re

from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Pt, Inches as PptxInches

from .layer_service import LayerService
from .message_service import MessageService


class ReportService:

    def __init__(self):
        self.layerService = LayerService()
        self.messageService = MessageService()

    @staticmethod
    def paragraphReplaceText(paragraph, regex, replace_str):
        """
        Replace all matches for regex in a paragraph while preserving runs.
        """

        while True:
            text = paragraph.text
            match = regex.search(text)

            if not match:
                break

            runs = iter(paragraph.runs)
            start, end = match.start(), match.end()

            for run in runs:
                run_len = len(run.text)

                if start < run_len:
                    break

                start, end = start - run_len, end - run_len

            run_text = run.text
            run_len = len(run_text)
            run.text = "%s%s%s" % (
                run_text[:start],
                replace_str,
                run_text[end:]
            )

            end -= run_len

            for run in runs:
                if end <= 0:
                    break

                run_text = run.text
                run_len = len(run_text)
                run.text = run_text[end:]
                end -= run_len

        return paragraph

    @staticmethod
    def normalizeImagePath(imageData):
        """
        Normalize image data to a file path.

        Supported inputs:
            - string path
            - list or tuple where the first item is the path
            - None
        """

        if imageData is None:
            return None

        if isinstance(imageData, str):
            return imageData

        if isinstance(imageData, (list, tuple)):
            if len(imageData) == 0:
                return None

            return imageData[0]

        return None

    def addImageInParagraph(self, document, imageData, feedback):
        totalData = len(imageData)
        progressPerFeature = 100.0 / totalData if totalData else 0

        for placeholder, value in imageData.items():

            if feedback.isCanceled():
                self.messageService.criticalMessageBar(
                    'Exporting maps',
                    'operation aborted by the user!'
                )
                break

            imagePath = self.normalizeImagePath(value)
            imageWidth = value[1] if isinstance(value, (list, tuple)) and len(value) > 1 else 4.0

            if imagePath is None or not os.path.isfile(imagePath):
                continue

            regex = re.compile(placeholder)

            for index, paragraph in enumerate(document.paragraphs):
                if placeholder in paragraph.text:
                    run = paragraph.add_run()
                    run.add_picture(imagePath, width=Inches(imageWidth))
                    self.paragraphReplaceText(paragraph, regex, '')
                    feedback.setProgress(int(index * progressPerFeature))

        feedback.setProgress(100)

    def addImageInTable(self, document, imageData, feedback):
        totalData = len(imageData)
        progressPerFeature = 100.0 / totalData if totalData else 0

        for placeholder, value in imageData.items():

            if feedback.isCanceled():
                self.messageService.criticalMessageBar(
                    'Exporting maps',
                    'operation aborted by the user!'
                )
                break

            imagePath = self.normalizeImagePath(value)
            imageWidth = value[1] if isinstance(value, (list, tuple)) and len(value) > 1 else 4.0

            if imagePath is None or not os.path.isfile(imagePath):
                continue

            regex = re.compile(placeholder)

            for index, table in enumerate(document.tables):
                for row in table.rows:
                    for cell in row.cells:
                        for cellIndex, paragraph in enumerate(cell.paragraphs):
                            if placeholder in paragraph.text:
                                run = paragraph.add_run()
                                run.add_picture(imagePath, width=Inches(imageWidth))
                                self.paragraphReplaceText(paragraph, regex, '')
                                feedback.setProgress(int(index * progressPerFeature))

        feedback.setProgress(100)

    def fillPlaceholdersOnParagraphs(self, document, trialData, feedback):
        totalData = len(trialData)
        progressPerFeature = 100.0 / totalData if totalData else 0

        for field, value in trialData.items():

            if feedback.isCanceled():
                self.messageService.criticalMessageBar(
                    'Exporting maps',
                    'operation aborted by the user!'
                )
                break

            regex = re.compile(field)

            for index, paragraph in enumerate(document.paragraphs):
                self.paragraphReplaceText(paragraph, regex, str(value))
                feedback.setProgress(int(index * progressPerFeature))

        feedback.setProgress(100)

    def fillPlaceholdersOnTable(self, document, tableData, feedback):
        totalData = len(tableData)
        progressPerFeature = 100.0 / totalData if totalData else 0

        for cellField, value in tableData.items():

            if feedback.isCanceled():
                self.messageService.criticalMessageBar(
                    'Exporting maps',
                    'operation aborted by the user!'
                )
                break

            cellRegex = re.compile(cellField)

            for index, table in enumerate(document.tables):
                for row in table.rows:
                    for cell in row.cells:
                        for cellIndex, paragraph in enumerate(cell.paragraphs):
                            self.paragraphReplaceText(
                                paragraph,
                                cellRegex,
                                str(value)
                            )
                            feedback.setProgress(int(cellIndex * progressPerFeature))

        feedback.setProgress(100)

    def appendControlScaleNoteToWord(self, document, paragraphData):
        """
        Add the control-scale note at the end of the Word report.

        If the template already contains {CONTROL_SCALE_NOTE}, it will also
        be replaced by fillPlaceholdersOnParagraphs().
        """

        note = paragraphData.get('{CONTROL_SCALE_NOTE}', None)

        if note:
            paragraph = document.add_paragraph()
            run = paragraph.add_run(str(note))
            run.bold = True

    def createWordReport(self, paragraphData, tableData, imageData, filePath, feedback):
        doc = os.path.join(
            self.layerService.getReportPath(),
            'report_template.docx'
        )

        reportDocument = Document(doc)

        self.fillPlaceholdersOnParagraphs(reportDocument, paragraphData, feedback)
        self.fillPlaceholdersOnTable(reportDocument, tableData, feedback)
        self.addImageInParagraph(reportDocument, imageData, feedback)
        self.addImageInTable(reportDocument, imageData, feedback)
        self.appendControlScaleNoteToWord(reportDocument, paragraphData)

        output = os.path.join(filePath, 'output_report.docx')
        reportDocument.save(output)

    def createPresentation(self, presentationData, filePath):
        doc = os.path.join(
            self.layerService.getPresentationPath(),
            'presentation_template.pptx'
        )

        reportPresentation = Presentation(doc)

        controlScaleNote = presentationData.get(
            '_CONTROL_SCALE_NOTE',
            None
        )

        for slideIndex, slide in enumerate(reportPresentation.slides, start=1):

            for dataIndex, data in presentationData.items():

                if not isinstance(dataIndex, int):
                    continue

                if slideIndex == dataIndex:

                    for placeholderIndex, placeholderData in data.items():

                        if placeholderIndex == 1:
                            self.changeTextPlaceholder(
                                slide,
                                placeholderIndex,
                                placeholderData,
                                36
                            )

                        elif placeholderIndex == 15:
                            self.changeTextPlaceholder(
                                slide,
                                placeholderIndex,
                                placeholderData,
                                18
                            )

                        else:
                            imagePath = self.normalizeImagePath(placeholderData)

                            if imagePath is None:
                                continue

                            if not os.path.isfile(imagePath):
                                print(f"Invalid file path for image: {imagePath}")
                                continue

                            try:
                                placeholder = slide.placeholders[placeholderIndex]
                                placeholder.insert_picture(imagePath)

                            except Exception as exception:
                                print(f"Error inserting picture: {exception}")

            if controlScaleNote and slideIndex in [5, 6, 7]:
                self.addControlScaleNoteToSlide(slide, controlScaleNote)

        output = os.path.join(filePath, 'output_presentation.pptx')
        reportPresentation.save(output)

    def addControlScaleNoteToSlide(self, slide, text):
        """
        Add a small note to PowerPoint slides explaining the reference scale.

        The note is added only to the slides where final surfaces are displayed.
        """

        try:
            left = PptxInches(0.45)
            top = PptxInches(6.85)
            width = PptxInches(8.5)
            height = PptxInches(0.35)

            textbox = slide.shapes.add_textbox(left, top, width, height)
            textFrame = textbox.text_frame
            textFrame.clear()

            paragraph = textFrame.paragraphs[0]
            run = paragraph.add_run()
            run.text = text

            font = run.font
            font.name = 'Times New Roman'
            font.size = Pt(11)
            font.bold = True

        except Exception as exception:
            print(f"Error adding control scale note: {exception}")

    def changeTextPlaceholder(self, slide, placeholderIndex, text, size):
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == placeholderIndex:
                self.textFormatting(
                    shape,
                    text,
                    size,
                    bold=True,
                    italic=False
                )

    def textFormatting(self, shape, text, size, bold=False, italic=False):
        text_frame = shape.text_frame
        text_frame.clear()

        paragraph = text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = str(text)

        self.fontFormatting(
            run,
            size,
            bold=bold,
            italic=italic
        )

    @staticmethod
    def fontFormatting(run, size, bold=False, italic=False):
        font = run.font
        font.name = 'Times New Roman'
        font.size = Pt(size)
        font.bold = bold
        font.italic = italic

    def iterate_over_slides(self):
        doc = os.path.join(
            self.layerService.getPresentationPath(),
            'presentation_template.pptx'
        )

        reportPresentation = Presentation(doc)

        for i, slide in enumerate(reportPresentation.slides, start=1):
            print(f"Slide {i}")

            for shape in slide.shapes:
                print(f"  Shape: {shape.name}")

            if slide.shapes.title:
                print(f"  Title: {slide.shapes.title.text}")

            for placeholder in slide.placeholders:
                if placeholder.has_text_frame:
                    print(
                        f"  Placeholder {placeholder.placeholder_format.idx}: "
                        f"{placeholder.text}"
                    )
